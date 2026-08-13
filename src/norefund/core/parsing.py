"""Extract plain text from supported document formats."""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".py", ".json"}


def extract_text(path: Path) -> str:
    """Dispatch to the right parser based on file extension."""
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")  # noqa: E501
    parsers = {
        ".pdf": _read_pdf,
        ".pptx": _read_pptx,
        ".docx": _read_docx,
    }
    try:
        return parsers.get(ext, _read_text)(path)
    except Exception as exc:
        raise RuntimeError(f"Failed to parse '{path.name}': {exc}") from exc


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    texts = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    return "\n".join(texts)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(para.text for para in doc.paragraphs)


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")
